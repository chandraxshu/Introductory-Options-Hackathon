"""
META Option Analysis (safe expiry selection)

This script performs the full workflow required by your hackathon problem statement using META stock.
Changes from previous version:
- selects the first option expiry that is strictly *after* today (ensures T>0 and non-zero Greeks)
- defensive handling of market prices and implied vol calculation
- debug prints to surface key inputs

Run this script in an environment with internet access and the following packages installed:
  pip install yfinance numpy pandas scipy matplotlib python-pptx

USAGE
  python META_option_analysis.py

Notes
- This uses Black-Scholes (European) as a simple baseline. Real-world US equity options are American; Black-Scholes ignores early exercise and dividends. Use put-call parity or binomial methods for more precision.
- You can change TICKER, OPTION_TYPE, or choose a specific expiry/strike by editing the parameters below.
"""

import math
import datetime as dt
import numpy as np
import pandas as pd
import yfinance as yf
from scipy.stats import norm
from scipy.optimize import brentq
import matplotlib.pyplot as plt
from pptx import Presentation

# ---------------------------
# User-configurable settings
# ---------------------------
TICKER = "META"
OPTION_TYPE = "call"   # "call" or "put"
# If SPECIFIC_EXPIRY is None, script picks the first expiry strictly after today
SPECIFIC_EXPIRY = None   # e.g. '2026-02-20' or None
SPECIFIC_STRIKE = None   # e.g. 230.0 or None -> chooses ATM
RISK_FREE_RATE = 0.045  # default annual risk-free rate (update as needed)
HIST_WINDOW_DAYS = 252  # use ~252 trading days (~1 year)
SAVE_OUTPUTS = True
CONTRACT_MULTIPLIER = 100  # standard US equity options
MIN_T_DAYS = 1  # minimum 1 day to avoid zero-time edge cases

# ---------------------------
# Black-Scholes functions
# ---------------------------

def bs_price(S, K, T, r, sigma, option_type="call"):
    if T <= 0 or sigma <= 0:
        # immediate intrinsic value if expired or no vol
        if option_type == "call":
            return max(S - K, 0.0)
        else:
            return max(K - S, 0.0)
    d1 = (math.log(S / K) + (r + 0.5 * sigma ** 2) * T) / (sigma * math.sqrt(T))
    d2 = d1 - sigma * math.sqrt(T)
    if option_type == "call":
        return S * norm.cdf(d1) - K * math.exp(-r * T) * norm.cdf(d2)
    else:
        return K * math.exp(-r * T) * norm.cdf(-d2) - S * norm.cdf(-d1)


def bs_greeks(S, K, T, r, sigma, option_type="call"):
    # returns Greeks per share. Theta returned as $/day
    if T <= 0 or sigma <= 0:
        return {"delta": 0.0, "theta": 0.0, "vega": 0.0}
    d1 = (math.log(S / K) + (r + 0.5 * sigma ** 2) * T) / (sigma * math.sqrt(T))
    d2 = d1 - sigma * math.sqrt(T)
    pdf_d1 = math.exp(-0.5 * d1 ** 2) / math.sqrt(2 * math.pi)
    if option_type == "call":
        delta = norm.cdf(d1)
        theta = -(S * pdf_d1 * sigma) / (2 * math.sqrt(T)) - r * K * math.exp(-r * T) * norm.cdf(d2)
    else:
        delta = norm.cdf(d1) - 1
        theta = -(S * pdf_d1 * sigma) / (2 * math.sqrt(T)) + r * K * math.exp(-r * T) * norm.cdf(-d2)
    vega = S * pdf_d1 * math.sqrt(T)
    # Convert theta from per-year to per-day by dividing by trading days (~252)
    theta_per_day = theta / 252.0
    return {"delta": delta, "theta": theta_per_day, "vega": vega}


def implied_volatility(market_price, S, K, T, r, option_type="call", tol=1e-6, maxiter=100):
    # Invert Black-Scholes for implied vol using Brent's method
    if market_price is None or market_price <= 0 or T <= 0:
        return float('nan')
    def objective(sigma):
        return bs_price(S, K, T, r, sigma, option_type) - market_price
    low, high = 1e-6, 5.0
    try:
        iv = brentq(objective, low, high, maxiter=maxiter, xtol=tol)
        return iv
    except Exception:
        return float('nan')

# ---------------------------
# Fetch data
# ---------------------------
print("Fetching historical price data for", TICKER)
asset = yf.Ticker(TICKER)
end_date = dt.datetime.now()
start_date = end_date - dt.timedelta(days=HIST_WINDOW_DAYS * 1.5)
# get daily close prices
hist = asset.history(start=start_date.strftime('%Y-%m-%d'), end=end_date.strftime('%Y-%m-%d'), interval='1d')
if hist.empty:
    raise RuntimeError("Failed to fetch historical data. Make sure you have internet access and the ticker is valid.")

prices = hist['Close'].dropna()
S0 = float(prices.iloc[-1])
print(f"Current spot price: {S0:.2f}")

# historical volatility (annualized)
logrets = np.log(prices / prices.shift(1)).dropna()
daily_vol = logrets.std()
annualized_hist_vol = float(daily_vol * math.sqrt(252))
print(f"Historical (annualized) vol over {len(logrets)} daily returns: {annualized_hist_vol:.2%}")

# ---------------------------
# Select an option contract (first expiry strictly after today)
# ---------------------------
expiries = asset.options
if len(expiries) == 0:
    raise RuntimeError("No option chain available for this ticker via yfinance.")

selected_expiry = None
if SPECIFIC_EXPIRY is not None:
    # use user-specified expiry if provided
    if SPECIFIC_EXPIRY in expiries:
        selected_expiry = SPECIFIC_EXPIRY
    else:
        raise RuntimeError(f"Specified expiry {SPECIFIC_EXPIRY} not found in available expiries: {expiries[:5]}...")
else:
    today = dt.datetime.now().date()
    for e in expiries:
        exp_date = dt.datetime.strptime(e, '%Y-%m-%d').date()
        if exp_date > today:
            selected_expiry = e
            break
    if selected_expiry is None:
        # fallback to the farthest expiry
        selected_expiry = expiries[-1]

print("Selected expiry:", selected_expiry)
opt_chain = asset.option_chain(selected_expiry)
if OPTION_TYPE == 'call':
    options_df = opt_chain.calls.copy()
else:
    options_df = opt_chain.puts.copy()

# Choose strike
if SPECIFIC_STRIKE is not None:
    if SPECIFIC_STRIKE in options_df['strike'].values:
        options_df = options_df[options_df['strike'] == SPECIFIC_STRIKE]
    else:
        raise RuntimeError(f"Specified strike {SPECIFIC_STRIKE} not found for expiry {selected_expiry}.")

# Choose ATM strike (closest to spot) if multiple
options_df['strike_diff'] = (options_df['strike'] - S0).abs()
if options_df.empty:
    raise RuntimeError("No options rows found after filtering.\n")
atm_row = options_df.loc[options_df['strike_diff'].idxmin()]
K = float(atm_row['strike'])

# market price selection: prefer mid of bid/ask, fallback to lastPrice
market_bid = atm_row.get('bid', np.nan)
market_ask = atm_row.get('ask', np.nan)
market_last = atm_row.get('lastPrice', np.nan)
if not np.isnan(market_bid) and not np.isnan(market_ask) and (market_bid > 0 or market_ask > 0):
    market_mid = float((market_bid + market_ask) / 2.0)
elif not np.isnan(market_last) and market_last > 0:
    market_mid = float(market_last)
else:
    market_mid = float(np.nan)

print(f"Selected option: {OPTION_TYPE} | Strike={K} | Market mid price={market_mid}")

# compute time to expiry in years (use calendar days, then convert)
today = dt.datetime.now().date()
expiry_date = dt.datetime.strptime(selected_expiry, '%Y-%m-%d').date()
T_days = (expiry_date - today).days
# enforce minimum T_days to avoid zero
T_days = max(T_days, MIN_T_DAYS)
T_years = T_days / 365.0
print(f"Time to expiry: {T_days} calendar days ({T_years:.4f} years)")

# ---------------------------
# Price using Black-Scholes (with historical vol)
# ---------------------------
bs_price_hist_vol = bs_price(S0, K, T_years, RISK_FREE_RATE, annualized_hist_vol, OPTION_TYPE)
print(f"BS price using historical vol ({annualized_hist_vol:.2%}): {bs_price_hist_vol:.4f}")

# implied vol from market price (if available)
iv = implied_volatility(market_mid, S0, K, T_years, RISK_FREE_RATE, OPTION_TYPE)
print(f"Implied volatility (from market mid): {iv if not np.isnan(iv) else 'N/A'}")

# BS price using implied vol (sanity check)
used_sigma_for_pricing = iv if not np.isnan(iv) else annualized_hist_vol
bs_price_iv = bs_price(S0, K, T_years, RISK_FREE_RATE, used_sigma_for_pricing, OPTION_TYPE)
print(f"BS price using chosen vol ({used_sigma_for_pricing:.2%}): {bs_price_iv:.4f}")

# ---------------------------
# Greeks
# ---------------------------
greeks_hist = bs_greeks(S0, K, T_years, RISK_FREE_RATE, annualized_hist_vol, OPTION_TYPE)
greeks_iv = bs_greeks(S0, K, T_years, RISK_FREE_RATE, used_sigma_for_pricing, OPTION_TYPE)

print("\nGreeks using historical vol (per share):")
for g, val in greeks_hist.items():
    print(f"  {g.capitalize():6s}: {val:.6f}")

print("\nGreeks using chosen vol (per share):")
for g, val in greeks_iv.items():
    print(f"  {g.capitalize():6s}: {val:.6f}")

# ---------------------------
# Hedging suggestion
# ---------------------------
delta_contract = greeks_iv['delta'] * CONTRACT_MULTIPLIER
hedge_shares = -delta_contract  # short if delta positive for a long call
print(f"\nDelta hedge suggestion (per 1 contract): trade {hedge_shares:.1f} shares to be delta-neutral.")

# ---------------------------
# Simple diagnostics & plots
# ---------------------------
plt.figure(figsize=(10, 5))
prices.tail(180).plot(title=f"{TICKER} last 180 trading days")
plt.xlabel('Date')
plt.ylabel('Close price')
if SAVE_OUTPUTS:
    plt.savefig('meta_price_history.png', bbox_inches='tight')

# volatility comparison plot
plt.figure(figsize=(10, 4))
rolling_vol = logrets.rolling(window=21).std() * math.sqrt(252)
rolling_vol.plot(title='21-day rolling annualized vol (approx)')
plt.axhline(annualized_hist_vol, color='black', linestyle='--', label='12m hist vol')
plt.legend()
if SAVE_OUTPUTS:
    plt.savefig('meta_volatility.png', bbox_inches='tight')

# summary results dataframe
summary = pd.DataFrame({
    'metric': ['spot', 'strike', 'time_to_expiry_days', 'time_to_expiry_years', 'hist_vol', 'implied_vol', 'bs_price_hist_vol', 'bs_price_implied', 'market_mid'],
    'value': [S0, K, T_days, T_years, annualized_hist_vol, iv, bs_price_hist_vol, bs_price_iv, market_mid]
})

if SAVE_OUTPUTS:
    summary.to_csv('meta_option_summary.csv', index=False)

# ---------------------------
# Generate one-page memo (markdown)
# ---------------------------
memo_lines = [
    f"# One-page memo: META option analysis ({dt.date.today().isoformat()})\n",
    f"**Ticker:** {TICKER}\n",
    f"**Selected expiry:** {selected_expiry} ({T_days} days)\n",
    f"**Option type / Strike:** {OPTION_TYPE} / {K}\n",
    f"**Spot price:** {S0:.2f}\n",
    f"**Market mid price (option):** {market_mid if not np.isnan(market_mid) else 'N/A'}\n",
    f"**Black-Scholes price (historical vol):** {bs_price_hist_vol:.4f}\n",
    f"**Implied volatility (from market):** {iv if not np.isnan(iv) else 'N/A'}\n",
    f"**Historical vol (annualized):** {annualized_hist_vol:.2%}\n",
    "\n**Key interpretation:**\n",
    ("- If implied vol > historical vol: market is pricing higher future uncertainty (option appears relatively expensive vs hist vol).\n" if (not np.isnan(iv) and iv > annualized_hist_vol) else "- Implied vol <= historical vol: option is not expensive relative to realized vol.\n"),
    "\n**Primary risks:**\n",
    f"- Delta exposure: {greeks_iv['delta']:.4f} per share ({greeks_iv['delta']*CONTRACT_MULTIPLIER:.1f} per contract)\n",
    f"- Theta (time decay): {greeks_iv['theta']:.4f} $/day per share ({greeks_iv['theta']*CONTRACT_MULTIPLIER:.2f} $/day per contract)\n",
    f"- Vega (vol sensitivity): {greeks_iv['vega']:.4f} $ per 1 vol-point per share ({greeks_iv['vega']*CONTRACT_MULTIPLIER:.2f} per contract)\n",
    "\n**Basic hedge suggestion:**\n",
    f"- To neutralize delta for 1 long contract, trade {hedge_shares:.1f} shares (short if negative).\n",
    "\n**Monitoring suggestions:**\n",
    "- Monitor implied vol vs realized vol and significant corporate events (earnings, guidance).\n",
    "- Rebalance delta hedge daily (or as liquidity permits).\n",
    "\n**Caveats:** Black-Scholes is a simplified model and ignores early exercise and discrete dividends.\n"
]

memo_text = '\n'.join(memo_lines)
if SAVE_OUTPUTS:
    with open('meta_one_page_memo.md', 'w') as f:
        f.write(memo_text)

# ---------------------------
# Create a small slide deck (6 slides) with highlights
# ---------------------------
if SAVE_OUTPUTS:
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"META Option Analysis: {OPTION_TYPE.capitalize()} Strike {K}"
    slide.placeholders[1].text = f"Date: {dt.date.today().isoformat()}"

    # Slide 2: Data & Parameters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Data & Parameters'
    body = slide.shapes.placeholders[1].text_frame
    body.text = f"Spot: {S0:.2f}\nStrike: {K}\nExpiry: {selected_expiry} ({T_days} days)\nRisk-free rate: {RISK_FREE_RATE:.2%}\nHist vol: {annualized_hist_vol:.2%}\nImplied vol: {iv if not np.isnan(iv) else 'N/A'}"

    # Slide 3: Price comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Model vs Market Price'
    body = slide.shapes.placeholders[1].text_frame
    body.text = f"Market mid price: {market_mid if not np.isnan(market_mid) else 'N/A'}\nBS price (hist vol): {bs_price_hist_vol:.4f}\nBS price (implied vol): {bs_price_iv:.4f}"

    # Slide 4: Greeks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Key Greeks (per contract)'
    body = slide.shapes.placeholders[1].text_frame
    body.text = f"Delta: {greeks_iv['delta']*CONTRACT_MULTIPLIER:.2f} shares per contract\nTheta: {greeks_iv['theta']*CONTRACT_MULTIPLIER:.2f} $/day per contract\nVega: {greeks_iv['vega']*CONTRACT_MULTIPLIER:.2f} $ per 1 vol-pt per contract"

    # Slide 5: Hedging suggestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Hedging & Risks'
    body = slide.shapes.placeholders[1].text_frame
    body.text = f"Delta hedge: trade {hedge_shares:.1f} shares per long contract to be delta-neutral.\nMonitor Theta (time decay) and Vega (vol changes)."

    # Slide 6: Appendix / Files
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Files Output'
    body = slide.shapes.placeholders[1].text_frame
    body.text = 'Files saved: meta_price_history.png, meta_volatility.png, meta_option_summary.csv, meta_one_page_memo.md'

    prs.save('meta_option_deck.pptx')
    print('\nSaved slide deck: meta_option_deck.pptx and memo: meta_one_page_memo.md')

# Final summary printed to console
print('\nSummary table saved to meta_option_summary.csv (and printed below):')
print(summary.to_string(index=False))


print('\nAll done. Check generated files for charts, memo and deck.')
